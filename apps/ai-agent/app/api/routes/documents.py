"""
Document upload  →  S3  →  extract text  →  chunk  →  embed  →  Pinecone

Supported file types
--------------------
  .pdf          — pypdf  (text)
  .txt          — UTF-8 decode
  .docx         — python-docx
  .pptx         — python-pptx  (all slides)
  .png          — Bedrock Claude 3 vision  (OCR + description)
  .jpg / .jpeg  — Bedrock Claude 3 vision  (OCR + description)

S3 layout
---------
  documents/{user_id}/{timestamp}-{uuid}{ext}   ← original file
  processed/{user_id}/{doc_id}.txt              ← extracted text (stored after embedding)
"""

from fastapi import APIRouter, BackgroundTasks, HTTPException, UploadFile, File, Form
from typing import Optional, List
from pathlib import Path
import asyncio
import json
import logging
import io
import mimetypes
import re

from app.core.config import settings

router = APIRouter()
logger = logging.getLogger(__name__)

MAX_FILE_SIZE        = 20 * 1024 * 1024   # 20 MB
MAX_FILES_PER_UPLOAD = 10

ALLOWED_EXT = {".pdf", ".txt", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}

# Map image extension → MIME type used by Claude vision
_IMAGE_MIME = {
    ".png":  "image/png",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
}


# ── Service accessors ─────────────────────────────────────────────────────────

def _get_langchain_service():
    from app.services.langchain_service import langchain_service
    if langchain_service is None:
        raise HTTPException(503, detail="Embedding service unavailable. Check PINECONE_API_KEY.")
    return langchain_service


def _get_s3():
    """Return the S3Service singleton, or None if S3 is not configured."""
    if not settings.S3_BUCKET_NAME:
        return None
    try:
        from app.services.s3_service import get_s3_service
        return get_s3_service()
    except Exception as exc:
        logger.warning("S3 service unavailable: %s — uploads will not be backed up", exc)
        return None


# ── Text extractors ───────────────────────────────────────────────────────────

def _extract_pdf(content: bytes) -> str:
    """Extract text from a PDF using pypdf (already in requirements)."""
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(content: bytes) -> str:
    """Extract text from a DOCX file using python-docx."""
    import docx
    doc = docx.Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also pull text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n".join(paragraphs)


def _extract_pptx(content: bytes) -> str:
    """Extract text from all slides of a PPTX file using python-pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        if slide_lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
    return "\n\n".join(parts)


def _extract_image(content: bytes, media_type: str) -> str:
    """OCR / describe an image using Bedrock Claude 3 vision."""
    from app.services.bedrock_service import get_bedrock_service
    svc = get_bedrock_service()
    return svc.describe_image_bytes(content, media_type)


def _extract_text(content: bytes, filename: str) -> str:
    """Dispatch to the correct extractor based on file extension."""
    suffix = Path(filename).suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(content)
    elif suffix == ".txt":
        return content.decode("utf-8", errors="ignore")
    elif suffix == ".docx":
        return _extract_docx(content)
    elif suffix == ".pptx":
        return _extract_pptx(content)
    elif suffix in _IMAGE_MIME:
        return _extract_image(content, _IMAGE_MIME[suffix])

    # Fallback — try UTF-8
    return content.decode(errors="ignore")

# ── Knowledge-graph enrichment ───────────────────────────────────────────────

async def _extract_and_update_graph(text: str, user_id: str, filename: str) -> None:
    """
    Background task: ask Claude to extract concepts + relationships from
    the document text, then write them into the Neo4j knowledge graph.

    We send the first ~6 000 characters (enough context, cheap to call).
    Claude returns a JSON object like:
      {
        "concepts": [
          {"name": "...", "description": "...", "domain": "...", "difficulty": 1-5},
          ...
        ],
        "prerequisites": [
          {"concept": "...", "requires": "..."},
          ...
        ]
      }
    """
    try:
        from app.services.bedrock_service import get_bedrock_service
        from app.services import knowledge_graph_service as kg

        svc    = get_bedrock_service()
        sample = text[:6000].strip()

        prompt = (
            f"You are a knowledge-graph builder.\n"
            f"Analyse the following document excerpt and extract:\n"
            f"1. Key learning concepts (academic topics, terms, theories, methods).\n"
            f"2. Prerequisite relationships — which concept must be understood before another.\n\n"
            f"Return ONLY valid JSON in this exact schema (no markdown, no extra text):\n"
            f'{{\n'
            f'  "concepts": [{{"name": string, "description": string, "domain": string, "difficulty": 1-5}}],\n'
            f'  "prerequisites": [{{"concept": string, "requires": string}}]\n'
            f'}}\n\n'
            f"Document excerpt from '{filename}':\n{sample}"
        )

        raw = svc.generate(prompt=prompt)

        # Strip accidental markdown code fences
        raw = re.sub(r"```(?:json)?\s*", "", raw).strip().rstrip("`")

        data = json.loads(raw)
        concepts     = data.get("concepts", [])
        prerequisites = data.get("prerequisites", [])

        for c in concepts:
            name = (c.get("name") or "").strip()
            if not name:
                continue
            await kg.add_concept(
                name=name,
                description=c.get("description") or "",
                domain=c.get("domain") or "",
                difficulty=int(c.get("difficulty") or 1),
            )
            # Track which document introduced each concept
            await kg.link_document_concept(source=filename, user_id=user_id, concept_name=name)

        for rel in prerequisites:
            concept  = (rel.get("concept")  or "").strip()
            requires = (rel.get("requires") or "").strip()
            if concept and requires:
                await kg.link_prerequisite(concept=concept, prerequisite=requires)

        logger.info(
            "[GRAPH] '%s' → %d concepts, %d prerequisite edges added to Neo4j",
            filename, len(concepts), len(prerequisites),
        )

    except json.JSONDecodeError as exc:
        logger.warning("[GRAPH] Claude returned non-JSON for '%s': %s", filename, exc)
    except Exception as exc:
        logger.warning("[GRAPH] knowledge-graph update failed for '%s': %s", filename, exc)

# ── Route handlers ────────────────────────────────────────────────────────────

@router.get("/")
async def list_documents():
    """Return Pinecone index stats — document count per namespace."""
    try:
        svc = _get_langchain_service()
        index = svc.pc.Index(svc.index_name)
        stats = index.describe_index_stats()
        namespaces    = stats.get("namespaces", {}) if isinstance(stats, dict) else {}
        total_vectors = stats.get("total_vector_count", 0) if isinstance(stats, dict) else 0
        return {
            "status":        "ok",
            "index":         svc.index_name,
            "total_vectors": total_vectors,
            "namespaces":    list(namespaces.keys()),
        }
    except HTTPException:
        raise
    except Exception as exc:
        logger.error("Failed to describe Pinecone index: %s", exc)
        return {"status": "unavailable", "error": str(exc)}


@router.post("/upload")
async def upload_documents(
    background_tasks: BackgroundTasks,
    files:   Optional[List[UploadFile]] = File(None),
    user_id: Optional[str]             = Form(None),
):
    """
    Upload one or more files  →  S3 (backup)  →  extract text  →  embed  →  Pinecone.

    Returns
    -------
    JSON with count, namespace, s3_uploads list, and any per-file errors.
    """
    if not settings.ENABLE_PERSISTENCE:
        raise HTTPException(403, detail="Document upload is disabled (ENABLE_PERSISTENCE=false).")

    if not files:
        raise HTTPException(400, "No files uploaded.")

    if len(files) > MAX_FILES_PER_UPLOAD:
        raise HTTPException(400, f"Maximum {MAX_FILES_PER_UPLOAD} files per upload.")

    try:
        from langchain_core.documents import Document as LCDocument
    except ImportError:
        raise HTTPException(503, "LangChain core not installed.")

    lc_svc    = _get_langchain_service()
    s3_svc    = _get_s3()
    namespace = user_id or "anonymous"

    embedded_count: int        = 0
    s3_uploads:     List[dict] = []
    errors:         List[str]  = []

    for f in files:
        if not f.filename:
            continue

        suffix = Path(f.filename).suffix.lower()
        if suffix not in ALLOWED_EXT:
            raise HTTPException(
                400,
                f"Unsupported file type '{suffix}'. "
                f"Allowed: {', '.join(sorted(ALLOWED_EXT))}",
            )

        try:
            content = await f.read()
            if len(content) > MAX_FILE_SIZE:
                raise HTTPException(400, f"'{f.filename}' exceeds 20 MB limit.")

            # ── Step 1: back up raw bytes to S3 ──────────────────────────────
            s3_info: dict = {}
            if s3_svc:
                try:
                    s3_info = s3_svc.upload_fileobj(
                        file_obj     = io.BytesIO(content),
                        filename     = f.filename,
                        user_id      = namespace,
                        content_type = f.content_type or "application/octet-stream",
                    )
                    s3_uploads.append(s3_info)
                    logger.info(
                        "[S3] %s → s3://%s/%s",
                        f.filename, s3_info.get("bucket"), s3_info.get("s3_key"),
                    )
                except Exception as s3_exc:
                    logger.warning(
                        "[S3] Upload failed for '%s': %s  (processing continues)",
                        f.filename, s3_exc,
                    )

            # ── Step 2: extract plain text ────────────────────────────────────
            text = _extract_text(content, f.filename)
            if not text.strip():
                logger.warning("[UPLOAD] No text extracted from '%s' — skipping", f.filename)
                continue

            # ── Step 3: store extracted text back to S3 (processed/ prefix) ──
            if s3_svc and s3_info.get("doc_id"):
                try:
                    s3_svc.store_processed_text(namespace, s3_info["doc_id"], text)
                except Exception as exc:
                    logger.warning("[S3] Could not store processed text: %s", exc)

            # ── Step 4: build LangChain Document ─────────────────────────────
            meta = {
                "source":       f.filename,
                "user_id":      namespace,
                "content_type": f.content_type or "unknown",
                "file_type":    suffix.lstrip("."),
            }
            if s3_info.get("s3_key"):
                meta["s3_key"]    = s3_info["s3_key"]
                meta["s3_bucket"] = s3_info.get("bucket", "")

            doc = LCDocument(page_content=text, metadata=meta)

            # ── Step 5: chunk + upsert to Pinecone ───────────────────────────
            chunks = lc_svc.split_documents([doc])
            logger.info(
                "[UPLOAD] '%s' → %d chunks → Pinecone namespace='%s'",
                f.filename, len(chunks), namespace,
            )
            await lc_svc.upsert_documents(chunks, collection_name=namespace)
            embedded_count += 1

            # ── Step 6: update knowledge graph (background — non-blocking) ───
            background_tasks.add_task(
                _extract_and_update_graph, text, namespace, f.filename
            )

        except HTTPException:
            raise
        except Exception as exc:
            logger.error("[UPLOAD] Failed to process '%s': %s", f.filename, exc)
            errors.append(f"{f.filename}: {exc!s}")

    if errors and embedded_count == 0:
        raise HTTPException(500, f"All uploads failed: {'; '.join(errors)}")

    return {
        "count":      embedded_count,
        "namespace":  namespace,
        "message":    f"Successfully embedded {embedded_count} document(s).",
        "s3_uploads": s3_uploads,
        "errors":     errors,
    }


@router.post("/presign")
async def presign_upload(body: dict):
    """
    Generate a presigned PUT URL for direct browser-to-S3 upload.

    Body: { "user_id": str, "filename": str, "content_type": str (optional) }
    Returns: { "upload_url": str, "s3_key": str, "doc_id": str, "bucket": str }
    """
    s3_svc = _get_s3()
    if not s3_svc:
        raise HTTPException(503, detail="S3 is not configured (S3_BUCKET_NAME missing).")

    user_id      = (body.get("user_id") or "anonymous").strip()
    filename     = (body.get("filename") or "upload").strip()
    content_type = body.get("content_type") or mimetypes.guess_type(filename)[0] or "application/octet-stream"

    import uuid, os as _os
    from datetime import datetime as _dt
    doc_id = str(uuid.uuid4())
    ts     = _dt.utcnow().strftime("%Y%m%d%H%M%S")
    ext    = _os.path.splitext(filename)[1] or ""
    s3_key = f"documents/{user_id}/{ts}-{doc_id}{ext}"

    upload_url = s3_svc.get_presigned_upload_url(s3_key, content_type=content_type)
    return {
        "upload_url": upload_url,
        "s3_key":     s3_key,
        "doc_id":     doc_id,
        "bucket":     s3_svc.bucket,
        "content_type": content_type,
    }


@router.delete("/by-source")
async def delete_document_by_source(
    user_id: str,
    source: Optional[str] = None,
    s3_key: Optional[str] = None,
):
    """
    Remove a document's vectors from Pinecone and (optionally) its file from S3.

    Query params:
      user_id  — Pinecone namespace / S3 prefix
      source   — original filename (used as Pinecone metadata filter)
      s3_key   — S3 object key to delete
    """
    deleted_vectors = False
    deleted_s3      = False

    # ── 1. Remove vectors from Pinecone ───────────────────────────────────
    if source:
        try:
            lc_svc = _get_langchain_service()
            index  = lc_svc.pc.Index(lc_svc.index_name)
            index.delete(filter={"source": source}, namespace=user_id)
            deleted_vectors = True
            logger.info("[DELETE] Pinecone vectors removed: source=%s namespace=%s", source, user_id)
        except HTTPException:
            raise
        except Exception as exc:
            logger.warning("[DELETE] Could not remove Pinecone vectors for source='%s': %s", source, exc)

    # ── 2. Remove file from S3 ────────────────────────────────────────────
    if s3_key:
        s3_svc = _get_s3()
        if s3_svc:
            try:
                s3_svc.delete_file(s3_key)
                deleted_s3 = True
                logger.info("[DELETE] S3 object removed: %s", s3_key)
            except Exception as exc:
                logger.warning("[DELETE] Could not delete S3 object '%s': %s", s3_key, exc)

    # ── 3. Remove document node + orphaned mastery edges from knowledge graph
    graph_result: dict = {}
    if source:
        try:
            from app.services import knowledge_graph_service as kg
            graph_result = await kg.remove_document_from_graph(
                source=source, user_id=user_id
            )
        except Exception as exc:
            logger.warning("[DELETE] Graph cleanup failed for '%s': %s", source, exc)

    return {
        "success":         True,
        "deleted_vectors": deleted_vectors,
        "deleted_s3":      deleted_s3,
        "graph":           graph_result,
    }


@router.get("/download")
async def download_url(s3_key: str):
    """
    Generate a presigned GET URL for downloading a document from S3.

    Query param: s3_key (e.g. documents/user123/20260101-uuid.pdf)
    Returns: { "download_url": str, "expires_in": int }
    """
    s3_svc = _get_s3()
    if not s3_svc:
        raise HTTPException(503, detail="S3 is not configured.")

    if not s3_key:
        raise HTTPException(400, detail="s3_key is required.")

    if not s3_svc.file_exists(s3_key):
        raise HTTPException(404, detail="File not found in S3.")

    url = s3_svc.get_presigned_url(s3_key)
    return {"download_url": url, "s3_key": s3_key, "expires_in": 3600}

